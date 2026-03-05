#!/usr/bin/env python3
"""
Postprocesado de PPTX: busca textboxes con el marker [[EXPERIENCE_BLOCK]],
parsea las experiencias laborales (ROL | EMPRESA | FECHA + bullets) y las
reemplaza por bloques de texto separados (role, company, dates, bullets).

Uso:
  python split_experience_blocks.py <input.pptx> <output.pptx>

Si no hay ningún marker, copia el archivo al output sin cambios.
Salida por stderr para logs; exit 0 = éxito, exit 1 = error.
"""

from __future__ import annotations

import argparse
import math
import os
import re
import shutil
import sys
from pathlib import Path

DEBUG = os.environ.get("DEBUG") == "1"

# Dependencia: pip install python-pptx. Solo reportamos "falta pptx" si el módulo pptx no existe.
try:
    from pptx import Presentation
except ModuleNotFoundError as e:
    if getattr(e, "name", "") == "pptx":
        print("Error: no se pudo importar 'pptx'. Instalá con: python -m pip install python-pptx", file=sys.stderr)
        sys.exit(1)
    raise
except ImportError:
    print("Error: no se pudo importar 'pptx'. Instalá con: python -m pip install python-pptx", file=sys.stderr)
    sys.exit(1)

from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

MARKER = "[[EXPERIENCE_BLOCK]]"

# Layout: cabecera en 2 columnas (rol ~35%, empresa+fecha ~65%), bullets debajo a ancho completo
ROLE_WIDTH_RATIO = 0.35   # rol ocupa 35% del ancho
META_WIDTH_RATIO = 0.65   # empresa+fecha 65%
LINE_HEIGHT_RATIO = 1.35
GAP_RATIO = 0.45
GAP_AFTER_HEADER_RATIO = 0.35  # gap entre cabecera y bullets
GAP_BETWEEN_EXP_RATIO = 0.6    # gap entre experiencias

# 1 pt ≈ 12700 EMU en Drawing ML
PT_TO_EMU = 12700

# Estilos fijos cuando no hay referencia de contenido (no usar nada del marker)
STYLE_ROLE = {"size_pt": 16, "name": None, "color": None, "bold": True, "italic": False}
STYLE_COMPANY = {"size_pt": 12, "name": None, "color": None, "bold": False, "italic": False}
STYLE_DATE = {"size_pt": 11, "name": None, "color": None, "bold": False, "italic": True}
STYLE_BULLETS = {"size_pt": 11, "name": None, "color": None, "bold": False, "italic": False}
# Fallbacks para extracción por tipo (A1)
HEADER_STYLE_FALLBACK = {"size_pt": 12, "name": None, "color": None, "bold": False, "italic": False}
BULLET_STYLE_FALLBACK = {"size_pt": 11, "name": None, "color": None, "bold": False, "italic": False}


def pt_to_emu(pt_val):
    return int(round(float(pt_val) * PT_TO_EMU))


def _iter_all_shapes(shapes):
    """Recorre todos los shapes, incluyendo los que están dentro de grupos."""
    for s in shapes:
        yield s
        if hasattr(s, "shapes"):
            yield from _iter_all_shapes(s.shapes)


def get_shape_text(shape):
    """Extrae todo el texto plano del text frame. No usar runs para detectar el marker."""
    try:
        if not getattr(shape, "has_text_frame", False):
            return ""
        full = shape.text_frame.text
        return full if full is not None else ""
    except Exception:
        return ""


def shape_contains_marker(shape):
    """Detección robusta: marker en el texto completo del text frame."""
    full = get_shape_text(shape)
    return MARKER in full


def _is_bullet_line(s):
    """True si la línea es un bullet (empieza con •, -, *, ·)."""
    return s.strip().startswith(("•", "-", "*", "·"))


def _lookahead_has_bullets(lines, start_idx, count=3):
    """True si en las próximas `count` líneas no vacías hay al menos una bullet."""
    n = 0
    for j in range(start_idx, min(start_idx + count * 2, len(lines))):
        s = lines[j].strip()
        if not s:
            continue
        n += 1
        if n > count:
            return False
        if _is_bullet_line(s):
            return True
    return False


def _is_real_header_line(stripped, lines, i):
    """
    Un header debe tener "|" y cumplir al menos una condición:
    - >=2 pipes (formato "Rol | Empresa | Fecha"),
    - o estructura tipo "X | Y (YYYY-YYYY)" (un pipe con año al final),
    - o lookahead: las próximas líneas contienen bullets.
    Si no, se trata como body para no partir mal experiencias.
    """
    if "|" not in stripped or stripped.startswith(("•", "-", "*", "·")):
        return False
    pipe_count = stripped.count("|")
    if pipe_count >= 2:
        return True
    # Un solo pipe: aceptar si parece "ROL | EMPRESA (YYYY-YYYY)" (año 4 dígitos)
    if re.search(r"\(?\d{4}\s*[-–—]\s*\d{4}\)?|\d{4}", stripped):
        return True
    if _lookahead_has_bullets(lines, i + 1):
        return True
    return False


def parse_experiences(text):
    """
    Parsea texto de bloque de experiencia.
    Encabezado = línea que pasa _is_real_header_line (| con >=2 pipes, o año, o lookahead bullets).
    Bullets = líneas que empiezan con •, -, *, ·.
    Header: role | company | date; si 2 partes: role y date; si >3: primera=role, última=date, medio=company.
    """
    if not text or not text.strip():
        return []

    lines = [ln.rstrip() for ln in text.splitlines()]
    experiences = []
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue

        is_header = _is_real_header_line(stripped, lines, i)
        if is_header:
            parts = [p.strip() for p in stripped.split("|")]
            if len(parts) >= 3:
                role, company, dates = parts[0], parts[1], parts[2]
            elif len(parts) == 2:
                role, company, dates = parts[0], "", parts[1]
            else:
                role = parts[0] if parts else ""
                dates = parts[-1] if parts else ""
                company = " | ".join(parts[1:-1]) if len(parts) > 1 else ""
        else:
            # Línea sin header real: si ya hay una experiencia en curso, añadir como bullet; si no, una experiencia solo con "role" o body
            if experiences and not _is_bullet_line(stripped):
                experiences[-1]["bullets"].append(stripped)
                i += 1
                continue
            if experiences and _is_bullet_line(stripped):
                experiences[-1]["bullets"].append(stripped)
                i += 1
                continue
            role = stripped
            company = ""
            dates = ""

        bullets = []
        i += 1
        while i < len(lines):
            ln = lines[i]
            s = ln.strip()
            if not s:
                i += 1
                break
            if _is_bullet_line(s):
                bullets.append(s)
                i += 1
                continue
            if _is_real_header_line(s, lines, i):
                break
            bullets.append(s)
            i += 1

        experiences.append({
            "role": role or "",
            "company": company or "",
            "dates": dates or "",
            "bullets": bullets,
        })

    return experiences


def _get_paragraph_text_clean(para):
    """Texto del párrafo concatenando solo runs que no contienen el marker."""
    parts = []
    for run in para.runs:
        if MARKER in (run.text or ""):
            continue
        parts.append(run.text or "")
    return "".join(parts).strip()


def _extract_run_style(run):
    """Extrae size_pt, name, color, bold, italic del run (ignorar marker)."""
    info = {"size_pt": 11.0, "name": None, "color": None, "bold": False, "italic": False}
    try:
        if run.font.size is not None:
            info["size_pt"] = run.font.size.pt
    except Exception:
        pass
    if run.font.name:
        info["name"] = run.font.name
    if run.font.color and getattr(run.font.color, "type", None) is not None:
        try:
            rgb = run.font.color.rgb
            if rgb is not None:
                info["color"] = rgb
        except Exception:
            pass
    try:
        if run.font.bold is not None:
            info["bold"] = run.font.bold
    except Exception:
        pass
    try:
        if run.font.italic is not None:
            info["italic"] = run.font.italic
    except Exception:
        pass
    return info


def extract_experience_styles(shape):
    """
    Extrae header_style y bullet_style por separado del textbox original.
    - header_style: del primer párrafo "header real" (contiene "|", no es bullet) que no sea solo marker.
    - bullet_style: del primer párrafo que empiece con "•", "-" o "·" que no sea marker.
    Devuelve {"header_style": dict|None, "bullet_style": dict|None}. Si no se detecta, se usa fallback en process_slide.
    """
    result = {"header_style": None, "bullet_style": None}
    try:
        if not shape.has_text_frame:
            return result
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = _get_paragraph_text_clean(para)
            if not text:
                continue
            run_for_style = None
            for run in para.runs:
                if MARKER in (run.text or ""):
                    continue
                if (run.text or "").strip():
                    run_for_style = run
                    break
            if run_for_style is None:
                continue
            style = _extract_run_style(run_for_style)

            # Header: línea con "|" que no empiece con bullet
            if "|" in text and not text.startswith(("•", "-", "*", "·")):
                if result["header_style"] is None:
                    result["header_style"] = style
            # Bullet: primer párrafo que empiece con bullet
            if text.startswith(("•", "-", "·")) or (text.startswith("*") and len(text) > 1):
                if result["bullet_style"] is None:
                    result["bullet_style"] = style
            if result["header_style"] is not None and result["bullet_style"] is not None:
                break
        return result
    except Exception:
        return result


def remove_shape(shape):
    """Elimina el shape del slide (vía XML)."""
    try:
        parent = shape._element.getparent()
        if parent is not None:
            parent.remove(shape._element)
    except Exception:
        pass


def add_textbox(slide, left_emu, top_emu, width_emu, height_emu, text, style_info, bullet=False):
    """Añade un textbox con el estilo indicado (size_pt, name, color, bold, italic).
    No se reduce font size para 'hacer entrar' texto; preferir overflow antes que romper estilos."""
    try:
        left = Emu(left_emu)
        top = Emu(top_emu)
        width = Emu(width_emu)
        height = Emu(height_emu)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        size_pt = style_info.get("size_pt", 11)
        if size_pt is None or not isinstance(size_pt, (int, float)):
            size_pt = 11
        p.font.size = Pt(float(size_pt))
        if style_info.get("name"):
            p.font.name = style_info["name"]
        if style_info.get("color"):
            p.font.color.rgb = style_info["color"]
        p.font.bold = style_info.get("bold", False)
        p.font.italic = style_info.get("italic", False)
        if bullet:
            p.level = 0
        return box
    except Exception as e:
        print(f"Warning: no se pudo crear textbox: {e}", file=sys.stderr)
        return None


def _line_height_emu(size_pt):
    """Altura de una línea en EMU según tamaño de fuente (evita superposición)."""
    return pt_to_emu(size_pt * LINE_HEIGHT_RATIO)


def _gap_emu(size_pt):
    """Gap entre bloques en EMU según tamaño de fuente."""
    return pt_to_emu(size_pt * GAP_RATIO)


def _gap_after_header_emu(size_pt):
    """Gap entre cabecera (rol+meta) y bloque de bullets."""
    return pt_to_emu(size_pt * GAP_AFTER_HEADER_RATIO)


def _gap_between_exp_emu(size_pt):
    """Gap entre una experiencia y la siguiente."""
    return pt_to_emu(size_pt * GAP_BETWEEN_EXP_RATIO)


def _estimate_bullet_lines(bullet_text, width_emu, font_size_pt):
    """
    Estima cuántas líneas visuales ocupará el texto de bullets con el ancho y tamaño dados.
    Considera wrap: más caracteres o menos ancho => más líneas.
    """
    if not bullet_text or not bullet_text.strip():
        return 1
    total_chars = len(bullet_text)
    # Ancho en puntos (1 pt ≈ 12700 EMU)
    width_pt = width_emu / PT_TO_EMU
    # Caracteres por línea aproximados: ancho_pt / (font_size * factor típico ~0.55)
    chars_per_line = max(1, width_pt / (float(font_size_pt) * 0.55))
    lines_by_wrap = math.ceil(total_chars / chars_per_line)
    # Mínimo = líneas por saltos explícitos (cada bullet suele ser una línea)
    lines_by_newlines = bullet_text.count("\n") + 1
    return max(lines_by_newlines, lines_by_wrap)


def process_slide(slide, slide_idx, dry_run=False):
    """
    Procesa un slide: encuentra shapes con MARKER (incl. dentro de grupos), parsea, borra y recrea.
    El estilo se toma del primer contenido real (sin marker); el layout se calcula según el font size de cada bloque.
    """
    to_process = []
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if not getattr(shape, "has_text_frame", False):
                continue
            full = get_shape_text(shape)
            if MARKER not in full:
                continue
            full_clean = full.replace(MARKER, "").strip()
            experiences = parse_experiences(full_clean)
            if not experiences and full_clean:
                experiences = [{"role": "", "company": "", "dates": full_clean, "bullets": []}]
            styles = extract_experience_styles(shape)
            shape_id = getattr(shape, "shape_id", None)
            to_process.append((shape, experiences, styles, shape_id))
        except Exception as e:
            print(f"Warning: error inspeccionando shape: {e}", file=sys.stderr)

    for shape, experiences, styles, shape_id in to_process:
        print(f"[split_experience] slide={slide_idx + 1} marker_shape_id={shape_id} experiences={len(experiences)}")
        if dry_run:
            continue
        try:
            left_emu = shape.left
            top_emu = shape.top
            width_emu = shape.width
            height_emu = shape.height
            bottom_emu = top_emu + height_emu

            remove_shape(shape)

            header_style = styles.get("header_style") or HEADER_STYLE_FALLBACK
            bullet_style_base = styles.get("bullet_style") or BULLET_STYLE_FALLBACK
            role_style = {**header_style, "bold": True, "italic": False}
            meta_style = {**header_style, "bold": False, "italic": False}
            bullet_style = {**bullet_style_base, "bold": False, "italic": False}
            role_width_emu = int(width_emu * ROLE_WIDTH_RATIO)
            meta_width_emu = width_emu - role_width_emu
            meta_left_emu = left_emu + role_width_emu

            if DEBUG:
                print(f"[split_experience] DEBUG header_pt={header_style.get('size_pt')} bullet_pt={bullet_style_base.get('size_pt')}", file=sys.stderr)

            # Cursor vertical acumulativo: todo el área original (top_emu .. bottom_emu) es el presupuesto
            cursor_y = top_emu

            for exp in experiences:
                role_size = role_style["size_pt"]
                meta_size = meta_style["size_pt"]
                bullet_size = bullet_style["size_pt"]

                header_h = max(_line_height_emu(role_size), _line_height_emu(meta_size))
                gap_after_header_h = _gap_after_header_emu(max(role_size, meta_size))

                if exp["bullets"]:
                    bullet_text = "\n".join(exp["bullets"])
                    lines = 0
                    for b in exp["bullets"]:
                        lines += _estimate_bullet_lines(b, width_emu, bullet_size)
                    lines += 1  # padding
                    line_height_emu = _line_height_emu(bullet_size)
                    bullets_h = lines * line_height_emu
                    min_lines = max(1, len(exp["bullets"]) * 1)
                    bullets_h = max(bullets_h, min_lines * line_height_emu)
                    bullets_h += int(0.5 * line_height_emu)
                else:
                    bullet_text = ""
                    bullets_h = 0

                gap_between_h = _gap_between_exp_emu(bullet_size if exp["bullets"] else meta_size)

                # --- Fila 1: cabecera (mismo top = cursor_y) ---
                header_top = cursor_y
                if exp["role"]:
                    add_textbox(
                        slide, left_emu, header_top, role_width_emu, header_h,
                        exp["role"], role_style
                    )
                meta_parts = [exp["company"], exp["dates"]]
                meta_parts = [p for p in meta_parts if p]
                if meta_parts:
                    meta_text = f"{meta_parts[0]} ({meta_parts[1]})" if len(meta_parts) == 2 else meta_parts[0]
                    add_textbox(
                        slide, meta_left_emu, header_top, meta_width_emu, header_h,
                        meta_text, meta_style
                    )

                # --- Fila 2: bullets debajo, ancho completo; altura = estimada por contenido ---
                bullets_top = cursor_y + header_h + gap_after_header_h
                if exp["bullets"] and bullets_h > 0:
                    add_textbox(
                        slide, left_emu, bullets_top, width_emu, bullets_h,
                        bullet_text, bullet_style, bullet=True
                    )
                    cursor_y = bullets_top + bullets_h
                else:
                    cursor_y = bullets_top

                cursor_y += gap_between_h

        except Exception as e:
            print(f"Error procesando bloque de experiencia: {e}", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(description="Split experience blocks in PPTX (marker: [[EXPERIENCE_BLOCK]]).")
    parser.add_argument("input_pptx", type=Path, help="Ruta al PPTX de entrada")
    parser.add_argument("output_pptx", type=Path, help="Ruta al PPTX de salida")
    parser.add_argument("--dry-run", action="store_true", help="Solo detectar y loguear; no modificar el PPTX")
    args = parser.parse_args()

    input_path = args.input_pptx.resolve()
    output_path = args.output_pptx.resolve()

    if not input_path.exists():
        print(f"Error: el archivo de entrada no existe: {input_path}", file=sys.stderr)
        sys.exit(1)

    if input_path.suffix.lower() != ".pptx":
        print("Error: el archivo de entrada debe ser .pptx", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        print(f"Error: no se pudo abrir el PPTX (archivo corrupto o no válido): {e}", file=sys.stderr)
        sys.exit(1)

    # Detección robusta: usar text_frame.text y recorrer también shapes dentro de grupos
    found_blocks = 0
    for slide in prs.slides:
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if not getattr(shape, "has_text_frame", False):
                    continue
                full = get_shape_text(shape)
                if MARKER in full:
                    found_blocks += 1
            except Exception:
                continue

    if found_blocks == 0:
        print("[split_experience] found_blocks=0")
        try:
            shutil.copy2(str(input_path), str(output_path))
        except Exception as e:
            print(f"Error copiando archivo: {e}", file=sys.stderr)
            sys.exit(1)
        return

    print(f"[split_experience] found_blocks={found_blocks}")
    if args.dry_run:
        # Solo loguear en qué slide/shape hay marker y cuántas experiencias se detectarían
        for slide_idx, slide in enumerate(prs.slides):
            for shape in _iter_all_shapes(slide.shapes):
                try:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    full = get_shape_text(shape)
                    if MARKER not in full:
                        continue
                    full_clean = full.replace(MARKER, "").strip()
                    experiences = parse_experiences(full_clean)
                    if not experiences and full_clean:
                        experiences = [{"role": "", "company": "", "dates": full_clean, "bullets": []}]
                    sid = getattr(shape, "shape_id", None)
                    print(f"[split_experience] slide={slide_idx + 1} shapes=... marker_shape_id={sid} experiences={len(experiences)}")
                except Exception:
                    pass
        shutil.copy2(str(input_path), str(output_path))
        return

    for slide_idx, slide in enumerate(prs.slides):
        process_slide(slide, slide_idx, dry_run=False)

    try:
        prs.save(str(output_path))
    except Exception as e:
        print(f"Error guardando PPTX: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
