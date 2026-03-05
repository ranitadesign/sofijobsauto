#!/usr/bin/env python3
"""
Postprocesado de PPTX: busca textboxes con el marker [[SIDEBAR_BLOCK]],
parsea secciones del sidebar (título, underline, body) y las reemplaza por
textboxes independientes por sección.

Uso:
  python split_sidebar_blocks.py <input.pptx> <output.pptx> [--dry-run]

Si no hay ningún marker, copia input -> output sin cambios; exit 0.
"""

from __future__ import annotations

import argparse
import math
import os
import shutil
import sys
from pathlib import Path

DEBUG = os.environ.get("DEBUG") == "1"

try:
    from pptx import Presentation
except ModuleNotFoundError as e:
    if getattr(e, "name", "") == "pptx":
        print("Error: no se pudo importar 'pptx'. Instalá con: python -m pip install python-pptx", file=sys.stderr)
        sys.exit(1)
    raise
except ImportError:
    print("Error: no se pudo importar 'pptx'. Instalá con: python -m pip install python-pptx", file=sys.stderr)
    sys.exit(1)

from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN

MARKER = "[[SIDEBAR_BLOCK]]"
PT_TO_EMU = 12700
LINE_HEIGHT_RATIO = 1.35
GAP_BETWEEN_SECTIONS_RATIO = 0.5

# Preferir título distinto de CONTACTO para titleStyleRef (misma apariencia en todas las secciones)
CONTACTO_TITLE_NORMALIZED = "CONTACTO"

# Fallbacks cuando no se detecta el tipo (nunca usar title para body ni body para title)
STYLE_TITLE_FALLBACK = {"size_pt": 18, "name": None, "color": None, "bold": True, "italic": False, "alignment": None}
STYLE_LINE_FALLBACK = {"size_pt": 12, "name": None, "color": None, "bold": False, "italic": False, "alignment": None}
# Body fallback cuando no hay bullets en el sidebar (12pt normal)
STYLE_BODY_FALLBACK = {"size_pt": 12, "name": None, "color": None, "bold": False, "italic": False, "alignment": None}
# Límite máximo para title_size_pt (evitar locuras)
TITLE_SIZE_MAX_PT = 40
TITLE_BODY_RATIO = 2.0


def pt_to_emu(pt_val):
    return int(round(float(pt_val) * PT_TO_EMU))


def _iter_all_shapes(shapes):
    for s in shapes:
        yield s
        if hasattr(s, "shapes"):
            yield from _iter_all_shapes(s.shapes)


def get_shape_text(shape):
    try:
        if not getattr(shape, "has_text_frame", False):
            return ""
        return shape.text_frame.text or ""
    except Exception:
        return ""


def _is_underline_line(line):
    s = line.strip()
    if not s:
        return False
    return all(c in "─-_ \t" for c in s)


def _is_title_line(line):
    s = line.strip()
    if not s:
        return False
    if s.startswith(("•", "-", "*", "·")):
        return False
    if _is_underline_line(s):
        return False
    return True


def _is_bullet_paragraph(text):
    """True si el texto del párrafo empieza con bullet (•, -, ·, *)."""
    s = (text or "").strip()
    return s.startswith(("•", "-", "·")) or (s.startswith("*") and len(s) > 1)


def _normalize_section_title(title):
    """Normaliza título de sección para comparar (ej. CONTACTO)."""
    return (title or "").strip().upper()


def _is_bullet_dot_paragraph(text):
    """True si el párrafo empieza con bullet real "•"."""
    s = (text or "").strip()
    return s.startswith("•")


def _run_contains_marker(run):
    return MARKER in (run.text or "")


def _get_paragraph_text_clean(para):
    """Texto del párrafo concatenando solo runs que no contienen el marker."""
    parts = []
    for run in para.runs:
        if _run_contains_marker(run):
            continue
        t = run.text or ""
        parts.append(t)
    return "".join(parts).strip()


def _extract_run_style(run, para):
    """Extrae size_pt, name, color, bold, italic del run y alignment del párrafo."""
    info = {"size_pt": 11.0, "name": None, "color": None, "bold": False, "italic": False, "alignment": None}
    try:
        if run.font.size is not None:
            info["size_pt"] = run.font.size.pt
    except Exception:
        pass
    if run.font.name:
        info["name"] = run.font.name
    if run.font.color and getattr(run.font.color, "type", None) is not None:
        try:
            r = run.font.color.rgb
            if r is not None:
                info["color"] = r
        except Exception:
            pass
    try:
        if run.font.bold is not None:
            info["bold"] = run.font.bold
    except Exception:
        pass
    try:
        if run.font.italic is not None:
            info["italic"] = run.font.italic
    except Exception:
        pass
    try:
        if para.alignment is not None:
            info["alignment"] = para.alignment
    except Exception:
        pass
    return info


def extract_sidebar_styles(shape):
    """
    Extrae estilos de referencia globales UNA sola vez por bloque (ignorando el marker).
    bodyStyleRef se detecta de forma estricta desde el primer bullet "•" (o "-", "·" como fallback).
    Si no hay bullets en todo el sidebar, bodyStyleRef queda None (fallback 12pt en process_slide).

    - titleStyleRef: solo para name/color; el size del título se calcula como 2× body_size_pt.
    - lineStyleRef: primer párrafo underline; al aplicar se usa body size y bold=false.
    - bodyStyleRef: primer párrafo cuyo texto trim empiece con "•", o "-"/"·" como fallback.
      Extrae size_pt, font.name, color.rgb, italic, alignment del primer run no vacío.
    """
    result = {
        "titleStyleRef": None,
        "lineStyleRef": None,
        "bodyStyleRef": None,
    }
    first_title_style = None
    try:
        if not shape.has_text_frame:
            return result
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = _get_paragraph_text_clean(para)
            if not text:
                continue
            run_for_style = None
            for run in para.runs:
                if _run_contains_marker(run):
                    continue
                if (run.text or "").strip():
                    run_for_style = run
                    break
            if run_for_style is None:
                continue
            style = _extract_run_style(run_for_style, para)

            if _is_underline_line(text):
                if result["lineStyleRef"] is None:
                    result["lineStyleRef"] = style
            elif _is_title_line(text):
                if first_title_style is None:
                    first_title_style = style
                if _normalize_section_title(text) != CONTACTO_TITLE_NORMALIZED and result["titleStyleRef"] is None:
                    result["titleStyleRef"] = style
            elif _is_bullet_dot_paragraph(text):
                if result["bodyStyleRef"] is None:
                    result["bodyStyleRef"] = style
            elif text.strip().startswith(("-", "·")):
                if result["bodyStyleRef"] is None:
                    result["bodyStyleRef"] = style
        if result["titleStyleRef"] is None and first_title_style is not None:
            result["titleStyleRef"] = first_title_style
        return result
    except Exception:
        return result


def parse_sidebar_sections(text):
    """
    Parsea contenido del sidebar en secciones.
    Cada sección: title (línea que no es bullet ni underline), opcional underline, body_lines hasta el próximo title.
    """
    if not text or not text.strip():
        return []
    cleaned = text.replace(MARKER, "").strip()
    lines = [ln.rstrip() for ln in cleaned.splitlines()]
    sections = []
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        if not _is_title_line(line):
            i += 1
            continue
        title = stripped
        underline = ""
        i += 1
        if i < len(lines) and _is_underline_line(lines[i]):
            underline = lines[i].strip()
            i += 1
        body_lines = []
        while i < len(lines):
            ln = lines[i]
            s = ln.strip()
            if not s:
                i += 1
                continue
            if _is_title_line(ln) and not _is_underline_line(ln):
                break
            body_lines.append(ln.rstrip())
            i += 1
        sections.append({"title": title, "underline": underline, "body_lines": body_lines})
    return sections


def remove_shape(shape):
    try:
        parent = shape._element.getparent()
        if parent is not None:
            parent.remove(shape._element)
    except Exception:
        pass


def add_textbox(slide, left_emu, top_emu, width_emu, height_emu, text, style_info, bullet=False):
    try:
        box = slide.shapes.add_textbox(Emu(left_emu), Emu(top_emu), Emu(width_emu), Emu(height_emu))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text or ""
        size_pt = style_info.get("size_pt", 11)
        if size_pt is None or not isinstance(size_pt, (int, float)):
            size_pt = 11
        p.font.size = Pt(float(size_pt))
        if style_info.get("name"):
            p.font.name = style_info["name"]
        if style_info.get("color"):
            p.font.color.rgb = style_info["color"]
        p.font.bold = style_info.get("bold", False)
        p.font.italic = style_info.get("italic", False)
        if bullet:
            p.level = 0
        if style_info.get("alignment") is not None:
            try:
                p.alignment = style_info["alignment"]
            except Exception:
                pass
        return box
    except Exception as e:
        print(f"Warning: no se pudo crear textbox: {e}", file=sys.stderr)
        return None


def _line_height_emu(size_pt):
    return pt_to_emu(size_pt * LINE_HEIGHT_RATIO)


def _gap_between_sections_emu(size_pt):
    return pt_to_emu(size_pt * GAP_BETWEEN_SECTIONS_RATIO)


def _estimate_body_lines(body_text, width_emu, font_size_pt):
    if not body_text or not body_text.strip():
        return 1
    total_chars = len(body_text)
    width_pt = width_emu / PT_TO_EMU
    chars_per_line = max(1, width_pt / (float(font_size_pt) * 0.55))
    lines_by_wrap = math.ceil(total_chars / chars_per_line)
    lines_by_newlines = body_text.count("\n") + 1
    return max(lines_by_newlines, lines_by_wrap)


def process_slide(slide, slide_idx, dry_run=False):
    to_process = []
    for shape in _iter_all_shapes(slide.shapes):
        try:
            if not getattr(shape, "has_text_frame", False):
                continue
            full = get_shape_text(shape)
            if MARKER not in full:
                continue
            full_clean = full.replace(MARKER, "").strip()
            sections = parse_sidebar_sections(full_clean)
            if not sections and full_clean:
                sections = [{"title": "", "underline": "", "body_lines": [full_clean]}]
            styles = extract_sidebar_styles(shape)
            shape_id = getattr(shape, "shape_id", None)
            to_process.append((shape, sections, styles, shape_id))
        except Exception as e:
            print(f"Warning: error inspeccionando shape: {e}", file=sys.stderr)

    for shape, sections, styles, shape_id in to_process:
        print(f"[split_sidebar] slide={slide_idx + 1} marker_shape_id={shape_id} sections={len(sections)}")
        if dry_run:
            continue
        try:
            left_emu = shape.left
            top_emu = shape.top
            width_emu = shape.width
            height_emu = shape.height
            remove_shape(shape)

            body_style_ref = styles.get("bodyStyleRef") or STYLE_BODY_FALLBACK
            body_size_pt = float(body_style_ref.get("size_pt") or 12)
            body_size_pt = max(1, body_size_pt)
            title_size_pt = min(max(1, round(body_size_pt * TITLE_BODY_RATIO)), TITLE_SIZE_MAX_PT)
            title_style_ref = styles.get("titleStyleRef") or STYLE_TITLE_FALLBACK

            title_style = {
                "size_pt": title_size_pt,
                "name": body_style_ref.get("name") or title_style_ref.get("name"),
                "color": body_style_ref.get("color") or title_style_ref.get("color"),
                "bold": True,
                "italic": False,
                "alignment": None,
            }
            line_style_ref = styles.get("lineStyleRef") or STYLE_LINE_FALLBACK
            line_style = {
                "size_pt": body_size_pt,
                "name": body_style_ref.get("name") or line_style_ref.get("name"),
                "color": body_style_ref.get("color") or line_style_ref.get("color"),
                "bold": False,
                "italic": False,
                "alignment": line_style_ref.get("alignment"),
            }
            body_style_applied = {
                "size_pt": body_size_pt,
                "name": body_style_ref.get("name"),
                "color": body_style_ref.get("color"),
                "bold": False,
                "italic": body_style_ref.get("italic", False),
                "alignment": body_style_ref.get("alignment"),
            }
            cursor_y = top_emu

            if DEBUG:
                print(f"[split_sidebar] body_size_pt={body_size_pt}", file=sys.stderr)
                print(f"[split_sidebar] title_size_pt={title_size_pt} (2x)", file=sys.stderr)

            for sec in sections:
                is_contacto = _normalize_section_title(sec["title"]) == CONTACTO_TITLE_NORMALIZED
                title_pt = title_size_pt
                ul_pt = body_size_pt
                body_pt = body_size_pt

                title_h = _line_height_emu(title_pt)
                if sec["title"]:
                    add_textbox(slide, left_emu, cursor_y, width_emu, title_h, sec["title"], title_style)
                cursor_y += title_h

                if sec["underline"]:
                    ul_h = _line_height_emu(ul_pt)
                    add_textbox(slide, left_emu, cursor_y, width_emu, ul_h, sec["underline"], line_style)
                    cursor_y += ul_h

                body_pt_use = body_pt
                body_text = "\n".join(sec["body_lines"])
                if sec["body_lines"]:
                    use_bullet = not is_contacto
                    style_for_body = {**body_style_applied, "bold": False}
                    body_h = _estimate_body_lines(body_text, width_emu, body_pt_use) * _line_height_emu(body_pt_use)
                    add_textbox(slide, left_emu, cursor_y, width_emu, body_h, body_text, style_for_body, bullet=use_bullet)
                    cursor_y += body_h
                    if DEBUG:
                        print(f"[split_sidebar] section title={sec['title'][:30]!r}", file=sys.stderr)

                cursor_y += _gap_between_sections_emu(max(title_pt, body_pt_use))

        except Exception as e:
            print(f"Error procesando bloque sidebar: {e}", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(description="Split sidebar blocks in PPTX (marker: [[SIDEBAR_BLOCK]]).")
    parser.add_argument("input_pptx", type=Path, help="Ruta al PPTX de entrada")
    parser.add_argument("output_pptx", type=Path, help="Ruta al PPTX de salida")
    parser.add_argument("--dry-run", action="store_true", help="Solo detectar y loguear; no modificar el PPTX")
    args = parser.parse_args()

    input_path = args.input_pptx.resolve()
    output_path = args.output_pptx.resolve()

    if not input_path.exists():
        print(f"Error: el archivo de entrada no existe: {input_path}", file=sys.stderr)
        sys.exit(1)
    if input_path.suffix.lower() != ".pptx":
        print("Error: el archivo de entrada debe ser .pptx", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(str(input_path))
    except Exception as e:
        print(f"Error: no se pudo abrir el PPTX: {e}", file=sys.stderr)
        sys.exit(1)

    found_blocks = 0
    for slide in prs.slides:
        for shape in _iter_all_shapes(slide.shapes):
            try:
                if not getattr(shape, "has_text_frame", False):
                    continue
                if MARKER in (get_shape_text(shape) or ""):
                    found_blocks += 1
            except Exception:
                continue

    if found_blocks == 0:
        print("[split_sidebar] found_blocks=0")
        try:
            shutil.copy2(str(input_path), str(output_path))
        except Exception as e:
            print(f"Error copiando archivo: {e}", file=sys.stderr)
            sys.exit(1)
        return

    print(f"[split_sidebar] found_blocks={found_blocks}")
    if args.dry_run:
        for slide_idx, slide in enumerate(prs.slides):
            for shape in _iter_all_shapes(slide.shapes):
                try:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    full = get_shape_text(shape)
                    if MARKER not in full:
                        continue
                    sections = parse_sidebar_sections(full.replace(MARKER, "").strip())
                    sid = getattr(shape, "shape_id", None)
                    print(f"[split_sidebar] slide={slide_idx + 1} marker_shape_id={sid} sections={len(sections)}")
                except Exception:
                    pass
        shutil.copy2(str(input_path), str(output_path))
        return

    for slide_idx, slide in enumerate(prs.slides):
        process_slide(slide, slide_idx, dry_run=False)

    try:
        prs.save(str(output_path))
    except Exception as e:
        print(f"Error guardando PPTX: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
